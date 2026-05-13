from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "SecurePlan-AI-Presentation.pptx"

slides = [
    ("SecurePlan AI", "Agentic security test planning from a web application URL"),
    ("Problem", "Security teams spend significant time converting app behavior into clear, structured security test plans."),
    ("Solution", "SecurePlan AI explores the UI, maps evidence to OWASP risk areas, and generates Markdown, JSON, and text reports."),
    ("How It Works", "URL intake -> UI exploration -> OWASP mapping -> AI scenario generation -> report export"),
    ("Core Features", "Professional dashboard, agent timeline, discovered flows, OWASP mapping, scenario table, export center"),
    ("Architecture", "React + FastAPI + exploration agent + scanner enrichment + OpenAI synthesis + report builder"),
    ("Demo Flow", "Enter URL, choose passive mode, generate report, review scenarios, export artifacts"),
    ("Impact", "Faster planning, consistent security coverage, explainable evidence, reusable documentation"),
]


def add_slide(prs: Presentation, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(246, 247, 249)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.8), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_run = title_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(34)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(18, 24, 38)

    body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(8.4), Inches(3.0))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.text = body
    body_run = body_frame.paragraphs[0].runs[0]
    body_run.font.size = Pt(22)
    body_run.font.color.rgb = RGBColor(56, 68, 87)

    accent = slide.shapes.add_shape(1, Inches(0.75), Inches(5.25), Inches(8.6), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 180, 160)
    accent.line.color.rgb = RGBColor(0, 180, 160)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    for title, body in slides:
        add_slide(prs, title, body)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()

